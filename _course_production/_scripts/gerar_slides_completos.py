#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Script para gerar PPTX completo com todos os 12 slides da Sección 1
Usa python-pptx para criar apresentação profissional
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Configurações
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)  # 16:9 ratio

# Cores (em RGB)
COLORS = {
    "azul_principal": RGBColor(30, 64, 175),  # #1E40AF
    "azul_claro": RGBColor(59, 130, 246),  # #3B82F6
    "naranja": RGBColor(245, 158, 11),  # #F59E0B
    "verde": RGBColor(16, 185, 129),  # #10B981
    "gris_oscuro": RGBColor(55, 65, 81),  # #374151
    "blanco": RGBColor(255, 255, 255),  # #FFFFFF
    "gris_claro": RGBColor(240, 240, 240),  # #F0F0F0
    "rojo_claro": RGBColor(254, 226, 226),  # #FEE2E2
    "verde_claro": RGBColor(220, 252, 231),  # #DCFCE7
    "amarillo": RGBColor(254, 243, 199),  # #FEF3C7
}


def add_title_box(
    slide, x, y, width, height, text, font_size=40, color=None, bold=True
):
    """Helper para adicionar caixa de texto com formatação padrão"""
    if color is None:
        color = COLORS["azul_principal"]

    textbox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.text = text

    for paragraph in tf.paragraphs:
        paragraph.font.name = "Montserrat"
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color

    return textbox


def add_body_box(
    slide, x, y, width, height, text, font_size=20, color=None, align=PP_ALIGN.LEFT
):
    """Helper para adicionar caixa de corpo de texto"""
    if color is None:
        color = COLORS["gris_oscuro"]

    textbox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.text = text

    for paragraph in tf.paragraphs:
        paragraph.font.name = "Inter"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.alignment = align

    return textbox


def slide_1_portada(prs):
    """SLIDE 1: Portada"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Fundo gradiente
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = COLORS["azul_principal"]
    fill.gradient_stops[1].color.rgb = COLORS["azul_claro"]

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = "¿Blog Visual, Rápido y Gratis?"

    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = "Montserrat"
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS["blanco"]
    title_paragraph.alignment = PP_ALIGN.CENTER

    # Emojis
    emojis = ["💻", "🚀", "📝"]
    emoji_positions = [Inches(2.5), Inches(4.5), Inches(6.5)]

    for emoji, x_pos in zip(emojis, emoji_positions):
        emoji_box = slide.shapes.add_textbox(x_pos, Inches(3.8), Inches(1), Inches(0.8))
        emoji_frame = emoji_box.text_frame
        emoji_frame.text = emoji
        emoji_para = emoji_frame.paragraphs[0]
        emoji_para.font.size = Pt(36)
        emoji_para.font.color.rgb = COLORS["blanco"]
        emoji_para.alignment = PP_ALIGN.CENTER

    print("✅ Slide 1 (Portada)")


def slide_2_gancho(prs):
    """SLIDE 2: Gancho"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    # Columna izquierda
    add_title_box(slide, 0.4, 0.3, 4.5, 1, "¿Cansado de pagar\npor hosting?", 32)
    add_title_box(
        slide, 0.4, 1.2, 4.5, 1, "¿No quieres complicarte\ncon WordPress?", 32
    )

    # Buenas noticias
    bn_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(2.1), Inches(4.5), Inches(0.6)
    )
    bn_frame = bn_box.text_frame
    bn_frame.text = "✓ TENEMOS BUENAS NOTICIAS"
    bn_para = bn_frame.paragraphs[0]
    bn_para.font.name = "Poppins"
    bn_para.font.size = Pt(22)
    bn_para.font.bold = True
    bn_para.font.color.rgb = COLORS["naranja"]

    # Body
    body_text = "Puedes crear un blog\nPROFESIONAL\nGRATIS"
    body_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(2.8), Inches(4.5), Inches(1.5)
    )
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    lines = body_text.split("\n")
    body_frame.text = lines[0]

    for line in lines[1:]:
        p = body_frame.add_paragraph()
        p.text = line

    for paragraph in body_frame.paragraphs:
        paragraph.font.name = "Inter"
        paragraph.font.color.rgb = COLORS["gris_oscuro"]
        if "PROFESIONAL" in paragraph.text or "GRATIS" in paragraph.text:
            paragraph.font.bold = True
            paragraph.font.size = Pt(24)
        else:
            paragraph.font.size = Pt(18)

    # Imagen placeholder
    image_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(0.4), Inches(4.2), Inches(3.25)
    )
    image_shape.fill.solid()
    image_shape.fill.fore_color.rgb = COLORS["gris_claro"]
    image_shape.line.color.rgb = COLORS["verde"]
    image_shape.line.width = Pt(3)

    tf = image_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "[Tu blog screenshot aquí]\n(blog.dataengineer.net.br)"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.font.italic = True
    tf.vertical_anchor = 1

    print("✅ Slide 2 (Gancho)")


def slide_3_promesa_1(prs):
    """SLIDE 3: Promesa Punto 1"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    # Borde izquierdo
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(5.625)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = COLORS["naranja"]
    border.line.fill.background()

    add_title_box(slide, 0.5, 0.5, 9, 0.8, "1", 140, COLORS["azul_principal"], False)
    add_title_box(slide, 0.8, 1.2, 8, 1, "Instalar todo lo necesario", 40)
    add_body_box(
        slide,
        0.8,
        2.2,
        8,
        1.5,
        "Ruby • Git • Jekyll\nPaso a paso. Sin complicaciones.",
        24,
        COLORS["gris_oscuro"],
    )

    print("✅ Slide 3 (Promesa 1)")


def slide_4_promesa_2(prs):
    """SLIDE 4: Promesa Punto 2"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(5.625)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = COLORS["verde"]
    border.line.fill.background()

    add_title_box(slide, 0.5, 0.5, 9, 0.8, "2", 140, COLORS["verde"], False)
    add_title_box(slide, 0.8, 1.2, 8, 1, "Crear tu primer blog", 40)
    add_body_box(
        slide,
        0.8,
        2.2,
        8,
        1.5,
        "Con tus propios artículos\nEstructura profesional\nDiseño limpio",
        24,
        COLORS["gris_oscuro"],
    )

    print("✅ Slide 4 (Promesa 2)")


def slide_5_promesa_3(prs):
    """SLIDE 5: Promesa Punto 3"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(5.625)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = COLORS["azul_claro"]
    border.line.fill.background()

    add_title_box(slide, 0.5, 0.5, 9, 0.8, "3", 140, COLORS["azul_claro"], False)
    add_title_box(slide, 0.8, 1.2, 8, 1, "Publicar en GitHub Pages", 40)
    add_body_box(
        slide,
        0.8,
        2.2,
        8,
        1.5,
        "Tu dominio gratis\nOnline en minutos\nControl total",
        24,
        COLORS["gris_oscuro"],
    )

    print("✅ Slide 5 (Promesa 3)")


def slide_6_promesa_4(prs):
    """SLIDE 6: Promesa Punto 4"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(5.625)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = COLORS["naranja"]
    border.line.fill.background()

    add_title_box(slide, 0.5, 0.5, 9, 0.8, "4", 140, COLORS["naranja"], False)
    add_title_box(slide, 0.8, 1.2, 8, 1, "El siguiente paso", 40)
    add_body_box(
        slide,
        0.8,
        2.2,
        8,
        1.5,
        "Monetizar tu blog\nHacer crecer tu audiencia\nNuevas oportunidades",
        24,
        COLORS["gris_oscuro"],
    )

    print("✅ Slide 6 (Promesa 4)")


def slide_7_gran_promesa(prs):
    """SLIDE 7: Gran Promesa"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].color.rgb = COLORS["verde"]
    fill.gradient_stops[1].color.rgb = RGBColor(5, 150, 105)

    # Texto central
    add_body_box(
        slide,
        1,
        0.3,
        8,
        5.025,
        "CUANDO TERMINES ESTE CURSO\n\nTendrás un blog\nPROFESIONAL\n\n✓ Control total\n✓ Dominio gratis\n✓ Hosting gratis\n✓ NADA DE PAGAR\n\nNo una página de Wix.\nNo un blog de WordPress.\n\nUn blog REAL.\nQue TEMES que funcione.",
        28,
        COLORS["blanco"],
        PP_ALIGN.CENTER,
    )

    print("✅ Slide 7 (Gran Promesa)")


def slide_8_comparativa(prs):
    """SLIDE 8: Comparativa"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    # Título
    add_title_box(slide, 0.5, 0.3, 9, 0.6, "¿Por qué Jekyll y no WordPress?", 36)

    # Tabla comparativa
    # WordPress (izquierda)
    wp_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(4.5), Inches(3.8)
    )
    wp_box.fill.solid()
    wp_box.fill.fore_color.rgb = COLORS["rojo_claro"]
    wp_box.line.color.rgb = RGBColor(239, 68, 68)
    wp_box.line.width = Pt(2)

    tf = wp_box.text_frame
    tf.clear()
    wp_items = [
        "💰 WordPress",
        "$50-200 al año",
        "🐢 Lento (~3s)",
        "🔓 Vulnerable",
        "📦 Pesado",
    ]
    tf.text = wp_items[0]
    for item in wp_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    for paragraph in tf.paragraphs:
        paragraph.font.name = "Inter"
        paragraph.font.size = Pt(18)
        paragraph.font.bold = paragraph.text.startswith("💰")
        paragraph.alignment = PP_ALIGN.CENTER

    # Jekyll (direita)
    jekyll_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.2), Inches(4.5), Inches(3.8)
    )
    jekyll_box.fill.solid()
    jekyll_box.fill.fore_color.rgb = COLORS["verde_claro"]
    jekyll_box.line.color.rgb = COLORS["verde"]
    jekyll_box.line.width = Pt(2)

    tf = jekyll_box.text_frame
    tf.clear()
    jekyll_items = [
        "✨ Jekyll",
        "$0 Gratis",
        "⚡ Velocidad",
        "🔒 Seguro",
        "📝 Minimalista",
    ]
    tf.text = jekyll_items[0]
    for item in jekyll_items[1:]:
        p = tf.add_paragraph()
        p.text = item

    for paragraph in tf.paragraphs:
        paragraph.font.name = "Inter"
        paragraph.font.size = Pt(18)
        paragraph.font.bold = paragraph.text.startswith("✨")
        paragraph.alignment = PP_ALIGN.CENTER

    print("✅ Slide 8 (Comparativa)")


def slide_9_proof(prs):
    """SLIDE 9: Proof - Tu Blog"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    add_body_box(
        slide,
        0.5,
        0.3,
        9,
        0.6,
        "ESTO ES LO QUE LOGRAMOS CON JEKYLL",
        32,
        COLORS["azul_principal"],
        PP_ALIGN.CENTER,
    )

    # Placeholder para imagem
    image_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(7), Inches(3.5)
    )
    image_shape.fill.solid()
    image_shape.fill.fore_color.rgb = COLORS["gris_claro"]
    image_shape.line.color.rgb = COLORS["verde"]
    image_shape.line.width = Pt(2)

    tf = image_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "[Screenshot de tu blog en alta resolución]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(153, 153, 153)
    tf.vertical_anchor = 1

    # Texto overlay
    add_body_box(
        slide,
        0.5,
        4.8,
        9,
        0.7,
        "⚡ 95+ Lighthouse Score  |  Tu blog en 2 horas",
        20,
        COLORS["blanco"],
        PP_ALIGN.CENTER,
    )

    print("✅ Slide 9 (Proof)")


def slide_10_requisitos(prs):
    """SLIDE 10: Requisitos"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["blanco"]

    add_title_box(slide, 0.5, 0.3, 9, 0.6, "Antes de empezar", 36)

    # Requisitos
    reqs = [
        "✅ Windows 10 o superior",
        "✅ Conexión a Internet",
        "✅ ~30 minutos de tu tiempo",
        "✅ Voluntad de aprender",
    ]
    for i, req in enumerate(reqs):
        add_body_box(
            slide, 0.8, 1.2 + (i * 0.5), 8, 0.5, req, 22, COLORS["gris_oscuro"]
        )

    # Box de aviso
    warning_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.5), Inches(8.4), Inches(1.8)
    )
    warning_box.fill.solid()
    warning_box.fill.fore_color.rgb = COLORS["amarillo"]
    warning_box.line.color.rgb = COLORS["naranja"]
    warning_box.line.width = Pt(2)

    tf = warning_box.text_frame
    tf.word_wrap = True
    tf.text = "⚠️ AVISO IMPORTANTE\n\nEsto es fácil, pero requiere atención en cada paso.\nAfortunadamente, si te atascas, estoy aquí para ayudarte."

    for paragraph in tf.paragraphs:
        paragraph.font.name = "Inter"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(124, 45, 18)

    print("✅ Slide 10 (Requisitos)")


def slide_11_cta(prs):
    """SLIDE 11: Llamada a la Acción"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 0
    fill.gradient_stops[0].color.rgb = COLORS["blanco"]
    fill.gradient_stops[1].color.rgb = COLORS["gris_claro"]

    add_title_box(slide, 0.5, 0.8, 9, 0.8, "¿ESTÁS LISTO?", 44)
    add_body_box(
        slide,
        0.5,
        1.8,
        9,
        1,
        "En el siguiente video:\nInstalamos Ruby + Git + Jekyll",
        24,
        COLORS["gris_oscuro"],
        PP_ALIGN.CENTER,
    )

    # CTA Button (visual)
    cta_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3.2), Inches(3), Inches(0.7)
    )
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = COLORS["azul_principal"]
    cta_box.line.fill.background()

    tf = cta_box.text_frame
    tf.text = "👉 VER MÓDULO 1"
    p = tf.paragraphs[0]
    p.font.name = "Inter"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS["blanco"]
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1

    add_body_box(
        slide,
        0.5,
        4.1,
        9,
        0.8,
        "Deja un comentario abajo. Dime qué te emociona más.",
        16,
        RGBColor(107, 114, 128),
        PP_ALIGN.CENTER,
    )

    print("✅ Slide 11 (CTA)")


def slide_12_despedida(prs):
    """SLIDE 12: Despedida/Créditos"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].color.rgb = COLORS["azul_principal"]
    fill.gradient_stops[1].color.rgb = COLORS["azul_claro"]

    add_body_box(
        slide,
        0.5,
        0.8,
        9,
        0.8,
        "¡Nos vemos en el siguiente módulo!",
        40,
        COLORS["blanco"],
        PP_ALIGN.CENTER,
    )

    add_body_box(
        slide,
        0.5,
        1.8,
        9,
        1.2,
        "📚 GitHub Pages + Jekyll\n💻 Windows Setup\n🚀 Blog profesional sin costos",
        24,
        COLORS["blanco"],
        PP_ALIGN.CENTER,
    )

    # Datos personales
    add_body_box(
        slide,
        0.5,
        3.0,
        9,
        0.5,
        "Marcos Vasconcellos de Andrade",
        22,
        COLORS["blanco"],
        PP_ALIGN.CENTER,
    )
    add_body_box(
        slide,
        0.5,
        3.5,
        9,
        0.7,
        "Data Engineer | AI Specialist",
        14,
        RGBColor(219, 234, 254),
        PP_ALIGN.CENTER,
    )
    add_body_box(
        slide,
        0.5,
        4.1,
        9,
        0.5,
        "Más de 20 años de experiencia en tecnología e IA",
        13,
        RGBColor(219, 234, 254),
        PP_ALIGN.CENTER,
    )

    add_body_box(
        slide,
        0.5,
        4.7,
        9,
        0.7,
        "🔗 blog.dataengineer.net.br  |  📧 [tu-email]  |  💼 LinkedIn",
        12,
        COLORS["blanco"],
        PP_ALIGN.CENTER,
    )

    print("✅ Slide 12 (Despedida/Créditos)")


def main():
    """Función principal"""

    print("🎬 Gerando apresentação PowerPoint completa (12 slides)...")
    print("")

    # Cria apresentação
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Cria todos os slides
    slide_1_portada(prs)
    slide_2_gancho(prs)
    slide_3_promesa_1(prs)
    slide_4_promesa_2(prs)
    slide_5_promesa_3(prs)
    slide_6_promesa_4(prs)
    slide_7_gran_promesa(prs)
    slide_8_comparativa(prs)
    slide_9_proof(prs)
    slide_10_requisitos(prs)
    slide_11_cta(prs)
    slide_12_despedida(prs)

    # Salva arquivo
    output_path = r"c:\Users\vasco\git\blog\_video-course-docs\SECCION_1_COMPLETA.pptx"
    prs.save(output_path)

    print("")
    print("=" * 70)
    print("✅ SUCESSO! Apresentação completa criada:")
    print(f"📁 {output_path}")
    print("=" * 70)
    print("")
    print("📊 Resumo de slides criados:")
    print("  1. Portada/Intro (0:00-0:10)")
    print("  2. Gancho (0:10-0:45)")
    print("  3. Promesa Punto 1 (0:45-1:15)")
    print("  4. Promesa Punto 2 (1:15-1:40)")
    print("  5. Promesa Punto 3 (1:40-2:05)")
    print("  6. Promesa Punto 4 (2:05-2:30)")
    print("  7. Gran Promesa (2:30-3:10)")
    print("  8. Comparativa (3:10-3:50)")
    print("  9. Proof - Tu Blog (3:50-4:10)")
    print("  10. Requisitos (4:10-4:30)")
    print("  11. CTA (4:30-4:50)")
    print("  12. Despedida/Créditos (4:50-5:00)")
    print("")
    print("⭐ Agora:")
    print("  1. Abre el archivo en Google Slides")
    print("  2. Añade tu imagen de blog en Slide 9")
    print("  3. Ajusta email/LinkedIn en Slide 12")
    print("  4. Verifica colores y fonts")
    print("  5. ¡Listo para grabar!")
    print("")


if __name__ == "__main__":
    main()
